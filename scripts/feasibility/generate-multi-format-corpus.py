#!/usr/bin/env python3

from __future__ import annotations

import argparse
import datetime
import hashlib
import io
import json
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Callable

from lxml import etree as LET
from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.formatting.rule import ColorScaleRule
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt


FIXED_TIME = datetime.datetime(2026, 1, 1, tzinfo=datetime.timezone.utc)
ZIP_TIME = (2026, 1, 1, 0, 0, 0)
CORE_PROPERTY_NAMESPACES = {
    "cp": "http://schemas.openxmlformats.org/package/2006/metadata/core-properties",
    "dc": "http://purl.org/dc/elements/1.1/",
    "dcterms": "http://purl.org/dc/terms/",
    "dcmitype": "http://purl.org/dc/dcmitype/",
    "xsi": "http://www.w3.org/2001/XMLSchema-instance",
}
DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
OFFICE_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

for prefix, namespace in CORE_PROPERTY_NAMESPACES.items():
    ET.register_namespace(prefix, namespace)


def normalize_core_properties(data: bytes) -> bytes:
    root = ET.fromstring(data)
    for name in ("created", "modified"):
        element = root.find(f"dcterms:{name}", CORE_PROPERTY_NAMESPACES)
        if element is not None:
            element.text = "2026-01-01T00:00:00Z"
    return ET.tostring(root, encoding="utf-8", xml_declaration=False)


def normalize_zip_bytes(source_bytes: bytes) -> bytes:
    with zipfile.ZipFile(io.BytesIO(source_bytes), "r") as source:
        members = [(info.filename, source.read(info.filename)) for info in source.infolist()]

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as target:
        for name, data in sorted(members):
            if name == "docProps/core.xml":
                data = normalize_core_properties(data)
            elif Path(name).suffix.lower() in {".docx", ".xlsx", ".pptx"}:
                data = normalize_zip_bytes(data)
            info = zipfile.ZipInfo(name, ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            target.writestr(info, data)
    return buffer.getvalue()


def normalize_zip(path: Path) -> None:
    path.write_bytes(normalize_zip_bytes(path.read_bytes()))


def rewrite_zip(path: Path, transform: Callable[[dict[str, bytes]], None]) -> None:
    with zipfile.ZipFile(path, "r") as source:
        members = {info.filename: source.read(info.filename) for info in source.infolist()}
    transform(members)

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as target:
        for name, data in sorted(members.items()):
            target.writestr(name, data)
    path.write_bytes(normalize_zip_bytes(buffer.getvalue()))


def create_reference_image(path: Path) -> None:
    image = Image.new("RGB", (640, 360), "#f4f7fb")
    draw = ImageDraw.Draw(image)
    draw.rectangle((24, 24, 616, 336), fill="#ffffff", outline="#183b66", width=4)
    draw.rectangle((56, 72, 224, 240), fill="#176b87")
    draw.ellipse((288, 72, 456, 240), fill="#f0b429")
    draw.line((56, 288, 584, 288), fill="#d64550", width=8)
    path.parent.mkdir(parents=True, exist_ok=True)
    image.save(path, format="PNG", optimize=False)


def create_docx(path: Path, image_path: Path) -> None:
    document = Document()
    document.core_properties.title = "KDV DOCX representative"
    document.core_properties.author = "KDV feasibility"
    document.core_properties.created = FIXED_TIME
    document.core_properties.modified = FIXED_TIME

    section = document.sections[0]
    section.orientation = WD_ORIENT.LANDSCAPE
    section.page_width = Inches(11.69)
    section.page_height = Inches(8.27)
    section.top_margin = Inches(0.6)
    section.bottom_margin = Inches(0.6)
    section.left_margin = Inches(0.7)
    section.right_margin = Inches(0.7)

    header = section.header.paragraphs[0]
    header.text = "KDV multi-format feasibility"
    header.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    title = document.add_heading("Word layout reference", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.runs[0].font.color.rgb = RGBColor(24, 59, 102)

    paragraph = document.add_paragraph()
    paragraph.add_run("Mixed formatting: ").bold = True
    italic = paragraph.add_run("italic")
    italic.italic = True
    paragraph.add_run(", Japanese 日本語, emoji ✓, and a long line that exercises wrapping.")

    table = document.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    values = [
        ("Feature", "Value", "Status"),
        ("Typography", "18 pt / bold", "Ready"),
        ("Table", "Merged-independent", "Ready"),
        ("External link", "https://example.invalid/", "Blocked"),
    ]
    for row, row_values in zip(table.rows, values):
        for cell, value in zip(row.cells, row_values):
            cell.text = value
    for cell in table.rows[0].cells:
        for run in cell.paragraphs[0].runs:
            run.bold = True

    document.add_picture(str(image_path), width=Inches(4.2))
    caption = document.add_paragraph("Figure 1: deterministic geometry reference")
    caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in caption.runs:
        run.font.size = Pt(10)

    document.add_page_break()
    document.add_heading("Second page", level=1)
    document.add_paragraph(
        "The second page verifies page breaks, headers, margins, and navigation."
    )
    footer = section.footer.paragraphs[0]
    footer.text = "KDV fixture"
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER

    document.save(path)
    normalize_zip(path)


def create_external_image_docx(source: Path, path: Path) -> None:
    path.write_bytes(source.read_bytes())

    def transform(members: dict[str, bytes]) -> None:
        document = LET.fromstring(members["word/document.xml"])
        blip = document.find(f".//{{{DRAWINGML_NS}}}blip")
        if blip is None:
            raise RuntimeError("embedded image relationship was not found")
        embed_key = f"{{{OFFICE_REL_NS}}}embed"
        link_key = f"{{{OFFICE_REL_NS}}}link"
        relationship_id = blip.attrib.pop(embed_key)
        blip.attrib[link_key] = relationship_id
        members["word/document.xml"] = LET.tostring(document, encoding="utf-8")

        relationships = LET.fromstring(members["word/_rels/document.xml.rels"])
        relationship = relationships.find(
            f"./{{{PACKAGE_REL_NS}}}Relationship[@Id='{relationship_id}']"
        )
        if relationship is None:
            raise RuntimeError("image relationship was not found")
        relationship.attrib["Target"] = (
            "http://127.0.0.1:43127/kdv-external-reference.png"
        )
        relationship.attrib["TargetMode"] = "External"
        members["word/_rels/document.xml.rels"] = LET.tostring(
            relationships, encoding="utf-8"
        )

    rewrite_zip(path, transform)


def create_macro_marker_docx(source: Path, path: Path) -> None:
    path.write_bytes(source.read_bytes())

    def transform(members: dict[str, bytes]) -> None:
        members["word/vbaProject.bin"] = b"KDV-MACRO-MARKER"

    rewrite_zip(path, transform)


def create_oversized_docx(source: Path, path: Path) -> None:
    path.write_bytes(source.read_bytes())

    def transform(members: dict[str, bytes]) -> None:
        document = LET.fromstring(members["word/document.xml"])
        body = document.find(f".//{{{WORD_NS}}}body")
        if body is None:
            raise RuntimeError("document body was not found")
        paragraph = LET.Element(f"{{{WORD_NS}}}p")
        run = LET.SubElement(paragraph, f"{{{WORD_NS}}}r")
        text = LET.SubElement(run, f"{{{WORD_NS}}}t")
        text.text = "A" * (32 * 1024 * 1024)
        body.insert(max(len(body) - 1, 0), paragraph)
        members["word/document.xml"] = LET.tostring(document, encoding="utf-8")

    rewrite_zip(path, transform)


def write_manifest(output: Path) -> None:
    entries = []
    for path in sorted(output.iterdir()):
        if not path.is_file() or path.name in {"README.md", "manifest.json"}:
            continue
        entry: dict[str, object] = {
            "name": path.name,
            "bytes": path.stat().st_size,
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        }
        if path.suffix.lower() in {".docx", ".xlsx", ".pptx"}:
            with zipfile.ZipFile(path, "r") as archive:
                entry["zip_entries"] = len(archive.infolist())
                entry["zip_uncompressed_bytes"] = sum(
                    member.file_size for member in archive.infolist()
                )
        entries.append(entry)
    (output / "manifest.json").write_text(
        json.dumps({"fixtures": entries}, indent=2) + "\n",
        encoding="utf-8",
    )


def create_xlsx(path: Path) -> None:
    workbook = Workbook()
    workbook.properties.title = "KDV XLSX representative"
    workbook.properties.creator = "KDV feasibility"
    workbook.properties.created = FIXED_TIME
    workbook.properties.modified = FIXED_TIME

    sheet = workbook.active
    sheet.title = "Dashboard"
    sheet.sheet_view.showGridLines = False
    sheet.freeze_panes = "A4"
    sheet.merge_cells("A1:F1")
    sheet["A1"] = "Quarterly performance"
    sheet["A1"].font = Font(size=22, bold=True, color="FFFFFF")
    sheet["A1"].fill = PatternFill("solid", fgColor="183B66")
    sheet["A1"].alignment = Alignment(horizontal="center", vertical="center")
    sheet.row_dimensions[1].height = 34

    headers = ["Region", "Q1", "Q2", "Q3", "Total", "Target"]
    for column, value in enumerate(headers, 1):
        cell = sheet.cell(3, column, value)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="176B87")
        cell.alignment = Alignment(horizontal="center")

    data = [
        ("North", 120, 132, 148, "=SUM(B4:D4)", 390),
        ("South", 98, 121, 137, "=SUM(B5:D5)", 350),
        ("East", 143, 151, 166, "=SUM(B6:D6)", 430),
        ("West", 88, 105, 119, "=SUM(B7:D7)", 330),
    ]
    thin = Side(style="thin", color="B7C4CE")
    for row_index, row_values in enumerate(data, 4):
        for column, value in enumerate(row_values, 1):
            cell = sheet.cell(row_index, column, value)
            cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)
            if column > 1:
                cell.number_format = '#,##0'

    sheet.conditional_formatting.add(
        "E4:E7",
        ColorScaleRule(
            start_type="min",
            start_color="FCE8E6",
            mid_type="percentile",
            mid_value=50,
            mid_color="FFF4CC",
            end_type="max",
            end_color="D9EAD3",
        ),
    )
    for column, width in enumerate([18, 12, 12, 12, 14, 14], 1):
        sheet.column_dimensions[get_column_letter(column)].width = width

    chart = BarChart()
    chart.title = "Quarter totals"
    chart.y_axis.title = "Units"
    chart.x_axis.title = "Region"
    chart.add_data(Reference(sheet, min_col=5, min_row=3, max_row=7), titles_from_data=True)
    chart.set_categories(Reference(sheet, min_col=1, min_row=4, max_row=7))
    chart.height = 7
    chart.width = 13
    sheet.add_chart(chart, "A10")

    notes = workbook.create_sheet("Notes")
    notes["A1"] = "日本語"
    notes["A2"] = "Formula results require a trusted cached value; KDV does not recalculate."
    notes["A3"] = "Remote resources must not be fetched."
    notes.column_dimensions["A"].width = 72

    workbook.save(path)
    normalize_zip(path)


def create_stress_xlsx(path: Path, rows: int, columns: int) -> None:
    workbook = Workbook(write_only=True)
    sheet = workbook.create_sheet("Cells")
    sheet.append([f"C{column}" for column in range(1, columns + 1)])
    for row in range(1, rows + 1):
        sheet.append([row * column for column in range(1, columns + 1)])
    workbook.save(path)
    normalize_zip(path)


def create_pptx(path: Path, image_path: Path) -> None:
    presentation = Presentation()
    presentation.core_properties.title = "KDV PPTX representative"
    presentation.core_properties.author = "KDV feasibility"
    presentation.core_properties.created = FIXED_TIME
    presentation.core_properties.modified = FIXED_TIME
    presentation.slide_width = PptxInches(13.333)
    presentation.slide_height = PptxInches(7.5)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = PptxRGBColor(244, 247, 251)

    title_box = slide.shapes.add_textbox(
        PptxInches(0.8), PptxInches(0.5), PptxInches(11.8), PptxInches(0.8)
    )
    title = title_box.text_frame.paragraphs[0]
    title.text = "Presentation layout reference"
    title.font.size = PptxPt(28)
    title.font.bold = True
    title.font.color.rgb = PptxRGBColor(24, 59, 102)

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        PptxInches(0.9),
        PptxInches(1.7),
        PptxInches(2.1),
        PptxInches(2.1),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = PptxRGBColor(23, 107, 135)
    circle.line.color.rgb = PptxRGBColor(255, 255, 255)
    circle.text = "1"
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.paragraphs[0].font.size = PptxPt(40)
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(255, 255, 255)

    slide.shapes.add_picture(
        str(image_path),
        PptxInches(3.6),
        PptxInches(1.55),
        width=PptxInches(5.2),
    )

    text_box = slide.shapes.add_textbox(
        PptxInches(9.2), PptxInches(1.7), PptxInches(3.2), PptxInches(3.2)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.text = "Layout features"
    paragraph.font.size = PptxPt(22)
    paragraph.font.bold = True
    for text in ["Absolute positioning", "Embedded image", "Japanese 日本語"]:
        item = text_frame.add_paragraph()
        item.text = text
        item.level = 0
        item.font.size = PptxPt(18)

    slide2 = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = PptxRGBColor(255, 255, 255)
    heading = slide2.shapes.add_textbox(
        PptxInches(0.8), PptxInches(0.45), PptxInches(11.8), PptxInches(0.8)
    )
    heading.text_frame.paragraphs[0].text = "Chart and navigation"
    heading.text_frame.paragraphs[0].font.size = PptxPt(28)
    heading.text_frame.paragraphs[0].font.bold = True

    chart_data = ChartData()
    chart_data.categories = ["North", "South", "East", "West"]
    chart_data.add_series("Units", (400, 356, 460, 312))
    slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        PptxInches(1.2),
        PptxInches(1.5),
        PptxInches(7.1),
        PptxInches(4.8),
        chart_data,
    )
    note = slide2.shapes.add_textbox(
        PptxInches(8.8), PptxInches(2.0), PptxInches(3.4), PptxInches(2.2)
    )
    note.text_frame.paragraphs[0].text = "Animations and transitions are intentionally unsupported."
    note.text_frame.paragraphs[0].font.size = PptxPt(20)

    presentation.save(path)
    normalize_zip(path)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("output", type=Path)
    parser.add_argument("--stress-rows", type=int, default=5000)
    parser.add_argument("--stress-columns", type=int, default=20)
    args = parser.parse_args()

    args.output.mkdir(parents=True, exist_ok=True)
    image_path = args.output / "reference-shapes.png"
    create_reference_image(image_path)
    docx_path = args.output / "representative.docx"
    create_docx(docx_path, image_path)
    create_external_image_docx(docx_path, args.output / "external-image.docx")
    create_macro_marker_docx(docx_path, args.output / "macro-marker.docx")
    create_oversized_docx(docx_path, args.output / "oversized-document.docx")
    create_xlsx(args.output / "representative.xlsx")
    create_stress_xlsx(
        args.output / "stress-100k-cells.xlsx",
        args.stress_rows,
        args.stress_columns,
    )
    create_pptx(args.output / "representative.pptx", image_path)
    write_manifest(args.output)


if __name__ == "__main__":
    main()
